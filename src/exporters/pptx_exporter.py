"""
PowerPoint Exporter Module

Handles exporting presentations to PowerPoint (.pptx) format
using python-pptx library.
"""

from typing import List, Dict, Any, Optional
import logging
from pathlib import Path


class PowerPointExporter:
    """
    Exports presentations to PowerPoint format.
    """
    
    def __init__(self):
        """Initialize the PowerPointExporter."""
        self.logger = logging.getLogger(__name__)
    
    def export(
        self,
        slides: List[Dict[str, Any]],
        metadata: Dict[str, Any],
        output_path: str,
        theme: str = "default",
        **kwargs
    ) -> bool:
        """
        Export slides to PowerPoint format.
        
        Args:
            slides: List of slide dictionaries
            metadata: Presentation metadata
            output_path: Output file path
            theme: Theme to apply
            **kwargs: Additional parameters
            
        Returns:
            True if successful, False otherwise
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Set presentation properties
            if 'title' in metadata:
                prs.core_properties.title = metadata['title']
            if 'author' in metadata:
                prs.core_properties.author = metadata['author']
            
            # Add slides
            for slide_data in slides:
                self._add_slide_to_presentation(prs, slide_data, theme)
            
            # Ensure output directory exists
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            # Save presentation
            prs.save(output_path)
            self.logger.info(f"Successfully exported PowerPoint to {output_path}")
            return True
        
        except Exception as e:
            self.logger.error(f"Error exporting PowerPoint: {str(e)}")
            return False
    
    def _add_slide_to_presentation(
        self,
        prs,
        slide_data: Dict[str, Any],
        theme: str
    ):
        """
        Add a single slide to the presentation.
        
        Args:
            prs: Presentation object
            slide_data: Slide data dictionary
            theme: Theme name
        """
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # Choose layout
            slide_type = slide_data.get('type', 'content')
            if slide_type == 'title':
                slide = prs.slides.add_slide(prs.slide_layouts[0])
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Add title
            if 'title' in slide_data and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_data['title']
            
            # Add content
            content = slide_data.get('content', {})
            if isinstance(content, dict):
                # Add bullet points
                if 'bullet_points' in content and len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    for point in content['bullet_points']:
                        p = text_frame.add_paragraph()
                        p.text = str(point)
                        p.level = 0
        
        except Exception as e:
            self.logger.error(f"Error adding slide: {str(e)}")
    
    def _apply_theme(self, prs, theme: str):
        """Apply theme to presentation."""
        # Theme application logic here
        pass
    
    def validate_output_path(self, output_path: str) -> bool:
        """
        Validate output path for PowerPoint export.
        
        Args:
            output_path: Output file path
            
        Returns:
            True if valid, False otherwise
        """
        if not output_path.lower().endswith('.pptx'):
            self.logger.warning(f"Output path should end with .pptx: {output_path}")
        
        try:
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            return True
        except Exception as e:
            self.logger.error(f"Invalid output path: {str(e)}")
            return False
